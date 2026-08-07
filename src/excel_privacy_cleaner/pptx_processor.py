from __future__ import annotations

import csv
import hashlib
import hmac
import json
import os
import re
import secrets
import shutil
import tempfile
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .excel_processor import AliasBook, ProcessingOptions, replacement_for
from .ginza_japanese import GinzaEntityDetector, WORD_NLP_CONFIDENCE, WORD_NLP_DETECTION_RULE
from .presidio_japanese import JapanesePresidioDetector, entity_label


PPTX_SUPPORTED_EXTENSION = ".pptx"
PPTX_MACRO_EXTENSION = ".pptm"

PPTX_CANDIDATE_CATEGORIES = {"会社名", "氏名", "住所", "電話番号", "メールアドレス", "銀行名"}
# Detection rules below are an independent copy of word_processor.py's regex
# set, not a shared import -- see docs/pptx_anonymization_v1_design.md
# ("既存資産の再利用方針"): PPTX text is bullet/fragment-heavy compared to
# Word's flowing prose, so the two need to be free to diverge in tuning
# without risking Word's already holdout-evaluated behavior.
COMPANY_DESIGNATORS = ("株式会社", "有限会社", "合同会社", "医療法人", "学校法人", "社会福祉法人")
COMPANY_SUFFIX_DESIGNATORS = ("株式会社", "有限会社", "合同会社")
PPTX_NAME_CHARS = r"一-龯々〆ヵヶぁ-んァ-ヶーA-Za-z0-9"
PPTX_REVIEW_CONFIDENCE_THRESHOLD = 0.75
PPTX_AUDIT_SCHEMA = "pptx_audit_v1"


class UnsupportedPptxDocumentError(RuntimeError):
    pass


@dataclass(frozen=True)
class PptxTextRun:
    slide_index: int
    shape_id: int
    shape_name: str
    container_type: str
    paragraph_index: int
    run_index: int
    char_start: int
    char_end: int
    text: str
    paragraph_text: str
    table_row: int | None = None
    table_col: int | None = None
    placeholder_type: str | None = None

    @property
    def location_id(self) -> str:
        values = [
            str(self.slide_index),
            str(self.shape_id),
            self.container_type,
            str(self.table_row),
            str(self.table_col),
            str(self.paragraph_index),
            str(self.run_index),
            str(self.char_start),
            str(self.char_end),
        ]
        return ":".join(values)


@dataclass(frozen=True)
class PptxParagraphText:
    slide_index: int
    shape_id: int
    shape_name: str
    container_type: str
    paragraph_index: int
    text: str
    table_row: int | None = None
    table_col: int | None = None
    placeholder_type: str | None = None


@dataclass(frozen=True)
class PptxDocumentProperties:
    author: str
    last_modified_by: str
    title: str
    subject: str
    keywords: str
    category: str
    comments: str
    content_status: str
    identifier: str
    language: str
    version: str


@dataclass(frozen=True)
class UnsupportedPptxFeature:
    feature_type: str
    slide_index: int
    count: int
    detail: str = ""


@dataclass(frozen=True)
class PptxStructureInventory:
    source_path: Path
    source_sha256: str
    slide_count: int
    paragraphs: tuple[PptxParagraphText, ...]
    runs: tuple[PptxTextRun, ...]
    document_properties: PptxDocumentProperties
    unsupported_features: tuple[UnsupportedPptxFeature, ...]


@dataclass(frozen=True)
class PptxCandidate:
    candidate_id: str
    category: str
    slide_index: int
    shape_id: int
    shape_name: str
    container_type: str
    paragraph_index: int
    char_start: int
    char_end: int
    text: str
    detection_rule: str
    confidence: float
    source: str
    location_id: str
    affected_run_indices: tuple[int, ...] = ()
    table_row: int | None = None
    table_col: int | None = None
    placeholder_type: str | None = None
    original_category: str = ""
    property_name: str | None = None


@dataclass
class PptxReplacementDecision:
    candidate: PptxCandidate
    enabled: bool = True
    replacement: str = ""
    # True only when a reviewer looked at a review-required candidate and
    # explicitly decided to keep the original text, as distinct from simply
    # never having been reviewed yet -- mirrors WordReplacementDecision.
    excluded: bool = False


@dataclass(frozen=True)
class PptxConversionResult:
    output_path: Path
    warnings: tuple[str, ...]
    converted_run_count: int
    converted_property_count: int
    review_required_count: int
    csv_path: Path
    report_path: Path


PPTX_CATEGORY_ALIAS_KIND = {
    "会社名": "company",
    "氏名": "name",
    "住所": "address",
    "電話番号": "phone",
    "メールアドレス": "email",
}


def default_replacement_for_candidate(candidate: PptxCandidate, alias_book: AliasBook, options: ProcessingOptions) -> str:
    if candidate.category == "銀行名":
        return alias_book.get("financial_institution", candidate.text)
    kind = PPTX_CATEGORY_ALIAS_KIND.get(candidate.category)
    if kind:
        return replacement_for(kind, candidate.text, alias_book, options)
    return alias_book.get("text", candidate.text)


def candidate_requires_review(candidate: PptxCandidate) -> bool:
    return candidate.confidence < PPTX_REVIEW_CONFIDENCE_THRESHOLD


def build_default_decisions(
    candidates: tuple[PptxCandidate, ...],
    alias_book: AliasBook,
    options: ProcessingOptions,
) -> list[PptxReplacementDecision]:
    decisions: list[PptxReplacementDecision] = []
    for candidate in candidates:
        enabled = not candidate_requires_review(candidate)
        replacement = default_replacement_for_candidate(candidate, alias_book, options)
        decisions.append(PptxReplacementDecision(candidate=candidate, enabled=enabled, replacement=replacement))
    return decisions


def extract_pptx_structure(path: Path) -> PptxStructureInventory:
    validate_pptx_input(path)
    source_sha256 = _file_sha256(path)
    unsupported_features = inspect_pptx_package(path)

    presentation = Presentation(path)
    paragraphs: list[PptxParagraphText] = []
    runs: list[PptxTextRun] = []

    for slide_index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            _append_shape(paragraphs, runs, shape, slide_index=slide_index)
        if slide.has_notes_slide:
            # has_notes_slide is checked first because accessing
            # .notes_slide on a slide that doesn't have one creates an empty
            # notes slide as a side effect -- undesirable during a read-only
            # scan.
            notes_frame = slide.notes_slide.notes_text_frame
            _append_text_frame(
                paragraphs,
                runs,
                notes_frame,
                slide_index=slide_index,
                shape_id=-1,
                shape_name="スピーカーノート",
                container_type="notes",
                placeholder_type=None,
            )

    return PptxStructureInventory(
        source_path=path,
        source_sha256=source_sha256,
        slide_count=len(presentation.slides),
        paragraphs=tuple(paragraphs),
        runs=tuple(runs),
        document_properties=_document_properties(presentation),
        unsupported_features=unsupported_features,
    )


def _append_shape(
    paragraphs: list[PptxParagraphText],
    runs: list[PptxTextRun],
    shape: Any,
    *,
    slide_index: int,
) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _append_shape(paragraphs, runs, child, slide_index=slide_index)
        return

    placeholder_type = None
    if shape.is_placeholder:
        placeholder_type = str(shape.placeholder_format.type)

    if getattr(shape, "has_table", False):
        table = shape.table
        for row_index, row in enumerate(table.rows):
            for col_index, cell in enumerate(row.cells):
                _append_text_frame(
                    paragraphs,
                    runs,
                    cell.text_frame,
                    slide_index=slide_index,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    container_type="table_cell",
                    placeholder_type=None,
                    table_row=row_index,
                    table_col=col_index,
                )
        return

    if getattr(shape, "has_text_frame", False):
        _append_text_frame(
            paragraphs,
            runs,
            shape.text_frame,
            slide_index=slide_index,
            shape_id=shape.shape_id,
            shape_name=shape.name,
            container_type="shape_text",
            placeholder_type=placeholder_type,
        )


def _append_text_frame(
    paragraphs: list[PptxParagraphText],
    runs: list[PptxTextRun],
    text_frame: Any,
    *,
    slide_index: int,
    shape_id: int,
    shape_name: str,
    container_type: str,
    placeholder_type: str | None,
    table_row: int | None = None,
    table_col: int | None = None,
) -> None:
    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        paragraph_text = "".join(run.text for run in paragraph.runs)
        paragraphs.append(
            PptxParagraphText(
                slide_index=slide_index,
                shape_id=shape_id,
                shape_name=shape_name,
                container_type=container_type,
                paragraph_index=paragraph_index,
                text=paragraph_text,
                table_row=table_row,
                table_col=table_col,
                placeholder_type=placeholder_type,
            )
        )
        cursor = 0
        for run_index, run in enumerate(paragraph.runs):
            text = run.text or ""
            if text:
                runs.append(
                    PptxTextRun(
                        slide_index=slide_index,
                        shape_id=shape_id,
                        shape_name=shape_name,
                        container_type=container_type,
                        paragraph_index=paragraph_index,
                        run_index=run_index,
                        char_start=cursor,
                        char_end=cursor + len(text),
                        text=text,
                        paragraph_text=paragraph_text,
                        table_row=table_row,
                        table_col=table_col,
                        placeholder_type=placeholder_type,
                    )
                )
            cursor += len(text)


def _document_properties(presentation: Any) -> PptxDocumentProperties:
    props = presentation.core_properties
    return PptxDocumentProperties(
        author=str(props.author or ""),
        last_modified_by=str(props.last_modified_by or ""),
        title=str(props.title or ""),
        subject=str(props.subject or ""),
        keywords=str(props.keywords or ""),
        category=str(props.category or ""),
        comments=str(props.comments or ""),
        content_status=str(props.content_status or ""),
        identifier=str(props.identifier or ""),
        language=str(props.language or ""),
        version=str(props.version or ""),
    )


def _document_property_values(properties: PptxDocumentProperties) -> tuple[tuple[str, str], ...]:
    return (
        ("author", properties.author),
        ("last_modified_by", properties.last_modified_by),
        ("title", properties.title),
        ("subject", properties.subject),
        ("keywords", properties.keywords),
        ("category", properties.category),
        ("comments", properties.comments),
        ("content_status", properties.content_status),
        ("identifier", properties.identifier),
        ("language", properties.language),
        ("version", properties.version),
    )


def detect_pptx_candidates(path: Path) -> tuple[PptxStructureInventory, tuple[PptxCandidate, ...]]:
    inventory = extract_pptx_structure(path)
    return inventory, candidates_for_inventory(inventory)


def candidates_for_inventory(inventory: PptxStructureInventory) -> tuple[PptxCandidate, ...]:
    detector = JapanesePresidioDetector()
    ginza_detector = GinzaEntityDetector()
    candidates: list[PptxCandidate] = []
    seen: set[tuple[str, str, int, int, str, str]] = set()
    runs_by_location = _runs_by_paragraph_location(inventory.runs)

    for paragraph in inventory.paragraphs:
        if not paragraph.text.strip():
            continue
        runs = runs_by_location.get(_paragraph_location_id(paragraph), ())
        _append_text_candidates(candidates, seen, inventory, paragraph, paragraph.text, runs, detector, ginza_detector)

    for property_name, value in _document_property_values(inventory.document_properties):
        if value:
            _append_property_candidates(candidates, seen, inventory, property_name, value, detector, ginza_detector)

    return tuple(candidates)


class PptxPrivacyProcessor:
    def __init__(self) -> None:
        self.detector = JapanesePresidioDetector()
        self.alias_book = AliasBook()
        self.options = ProcessingOptions()
        self.inventory: PptxStructureInventory | None = None
        self.temp_dir: Path | None = None
        self.temp_pptx: Path | None = None

    def cleanup(self) -> None:
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir, ignore_errors=True)
        self.temp_dir = None
        self.temp_pptx = None
        self.inventory = None

    def scan(self, pptx_path: Path, options: ProcessingOptions | None = None) -> list[PptxReplacementDecision]:
        self.cleanup()
        self.options = options or ProcessingOptions()
        self.alias_book = AliasBook()
        self.temp_dir = Path(tempfile.mkdtemp(prefix="PptxPrivacyCleaner_"))
        self.temp_pptx = self.temp_dir / pptx_path.name
        shutil.copy2(pptx_path, self.temp_pptx)
        self.inventory = extract_pptx_structure(self.temp_pptx)
        candidates = candidates_for_inventory(self.inventory)
        return build_default_decisions(candidates, self.alias_book, self.options)

    def convert(
        self,
        source_path: Path,
        decisions: list[PptxReplacementDecision],
        output_dir: Path | None = None,
        write_artifacts: bool = True,
    ) -> PptxConversionResult:
        if not self.temp_pptx or not self.temp_pptx.exists() or self.inventory is None:
            raise RuntimeError("先に検査を実行してください。")

        warnings: list[str] = []
        enabled_decisions = [decision for decision in decisions if decision.enabled]

        review_required = [
            decision
            for decision in decisions
            if not decision.enabled and not decision.excluded and candidate_requires_review(decision.candidate)
        ]
        if review_required:
            preview = "、".join(f"{decision.candidate.category}:{decision.candidate.text}" for decision in review_required[:8])
            raise RuntimeError(
                f"未確認候補が {len(review_required)} 件あります。信頼度の低い検出のため内容を確認し、"
                f"問題なければ候補を有効にするか、誤りであれば手動修正のうえ再実行してください。対象例: {preview}"
            )

        if self.inventory.unsupported_features:
            feature_counts: dict[str, int] = {}
            for feature in self.inventory.unsupported_features:
                feature_counts[feature.feature_type] = feature_counts.get(feature.feature_type, 0) + feature.count
            detail = "、".join(f"{feature_type}: {count}件" for feature_type, count in sorted(feature_counts.items()))
            if not self.options.is_analysis:
                raise RuntimeError(f"未対応領域が検出されたため外部共有用の出力を停止しました。未対応領域: {detail}")
            warnings.append(
                f"未対応領域が検出されました(未対応領域: {detail})。未対応領域内の匿名化は保証しない。"
                "分析継続用の範囲で利用し、外部共有する場合は内容を必ず確認してください。"
            )

        paragraph_by_location = {
            _paragraph_location_id(paragraph): paragraph for paragraph in self.inventory.paragraphs
        }
        property_decisions = [decision for decision in enabled_decisions if decision.candidate.source == "document_property"]
        paragraph_decisions = [decision for decision in enabled_decisions if decision.candidate.source != "document_property"]
        for decision in paragraph_decisions:
            candidate = decision.candidate
            paragraph = paragraph_by_location.get(candidate.location_id)
            if paragraph is None or paragraph.text[candidate.char_start : candidate.char_end] != candidate.text:
                raise RuntimeError(
                    f"構造が変化しました。再スキャンしてください。対象候補: {candidate.category} 「{candidate.text}」"
                )

        decisions_by_location: dict[str, list[PptxReplacementDecision]] = {}
        for decision in paragraph_decisions:
            decisions_by_location.setdefault(decision.candidate.location_id, []).append(decision)
        for location_id, location_decisions in decisions_by_location.items():
            ordered = sorted(location_decisions, key=lambda item: item.candidate.char_start)
            for previous, current in zip(ordered, ordered[1:]):
                if previous.candidate.char_end > current.candidate.char_start:
                    raise RuntimeError(
                        "検出候補の範囲が重複しています: "
                        f"{previous.candidate.category} 「{previous.candidate.text}」 / "
                        f"{current.candidate.category} 「{current.candidate.text}」"
                    )

        decisions_by_property: dict[str, list[PptxReplacementDecision]] = {}
        for decision in property_decisions:
            property_name = decision.candidate.property_name or ""
            decisions_by_property.setdefault(property_name, []).append(decision)
        for property_name, property_decision_list in decisions_by_property.items():
            ordered = sorted(property_decision_list, key=lambda item: item.candidate.char_start)
            for previous, current in zip(ordered, ordered[1:]):
                if previous.candidate.char_end > current.candidate.char_start:
                    raise RuntimeError(f"文書プロパティ {property_name} の検出候補の範囲が重複しています。")

        presentation = Presentation(self.temp_pptx)
        live_paragraphs_by_location = _live_paragraphs_by_location(presentation)
        if len(live_paragraphs_by_location) != len(self.inventory.paragraphs):
            raise RuntimeError("構造が変化しました。再スキャンしてください。")

        converted_run_count = 0
        for location_id, location_decisions in decisions_by_location.items():
            paragraph_object = live_paragraphs_by_location.get(location_id)
            if paragraph_object is None:
                continue
            run_items = _paragraph_runs_with_offsets(paragraph_object)
            converted_run_count += _apply_paragraph_decisions(run_items, location_decisions)

        converted_property_count = 0
        for property_name, property_decision_list in decisions_by_property.items():
            current_value = getattr(presentation.core_properties, property_name, "") or ""
            new_value = current_value
            for start, end, replacement in sorted(
                (
                    (decision.candidate.char_start, decision.candidate.char_end, decision.replacement)
                    for decision in property_decision_list
                ),
                key=lambda item: item[0],
                reverse=True,
            ):
                new_value = new_value[:start] + replacement + new_value[end:]
            setattr(presentation.core_properties, property_name, new_value)
            converted_property_count += 1

        output_path = _make_pptx_output_path(source_path, output_dir=output_dir, options=self.options)
        presentation.save(output_path)

        residual = _find_residual_text(output_path, paragraph_decisions + property_decisions)
        if residual:
            detail = _format_residual_detail(residual)
            if not self.options.is_analysis:
                output_path.unlink(missing_ok=True)
                raise RuntimeError(
                    f"内部XML残存検査で匿名化対象の原文が検出されたため外部共有用の出力を停止しました。残存箇所: {detail}"
                )
            warnings.append(
                f"内部XML残存検査で匿名化対象の原文が検出されました(残存箇所: {detail})。"
                "本文以外の内部XMLに原文が残っている可能性があるため、外部共有する場合は内容を必ず確認してください。"
            )

        review_required_count = sum(1 for decision in decisions if candidate_requires_review(decision.candidate))
        csv_path = output_path.with_name(f"{output_path.stem}_検出変換結果.csv")
        report_path = output_path.with_name(f"{output_path.stem}_処理報告書.txt")
        result = PptxConversionResult(
            output_path=output_path,
            warnings=tuple(warnings),
            converted_run_count=converted_run_count,
            converted_property_count=converted_property_count,
            review_required_count=review_required_count,
            csv_path=csv_path,
            report_path=report_path,
        )
        if write_artifacts:
            audit_path = output_path.with_name(f"{output_path.stem}_監査報告.json")
            write_pptx_findings_csv(csv_path, decisions)
            write_pptx_processing_report(
                report_path,
                source_path=source_path,
                result=result,
                decisions=decisions,
                inventory=self.inventory,
                options=self.options,
            )
            write_pptx_audit_json(
                audit_path,
                source_path=source_path,
                result=result,
                decisions=decisions,
                inventory=self.inventory,
                options=self.options,
            )
        self.cleanup()
        return result


def _live_paragraphs_by_location(presentation: Any) -> dict[str, Any]:
    result: dict[str, Any] = {}
    for slide_index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            _collect_live_paragraphs(result, shape, slide_index=slide_index)
        if slide.has_notes_slide:
            _collect_text_frame_paragraphs(
                result,
                slide.notes_slide.notes_text_frame,
                slide_index=slide_index,
                shape_id=-1,
                container_type="notes",
            )
    return result


def _collect_live_paragraphs(result: dict[str, Any], shape: Any, *, slide_index: int) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _collect_live_paragraphs(result, child, slide_index=slide_index)
        return
    if getattr(shape, "has_table", False):
        table = shape.table
        for row_index, row in enumerate(table.rows):
            for col_index, cell in enumerate(row.cells):
                _collect_text_frame_paragraphs(
                    result,
                    cell.text_frame,
                    slide_index=slide_index,
                    shape_id=shape.shape_id,
                    container_type="table_cell",
                    table_row=row_index,
                    table_col=col_index,
                )
        return
    if getattr(shape, "has_text_frame", False):
        _collect_text_frame_paragraphs(
            result,
            shape.text_frame,
            slide_index=slide_index,
            shape_id=shape.shape_id,
            container_type="shape_text",
        )


def _collect_text_frame_paragraphs(
    result: dict[str, Any],
    text_frame: Any,
    *,
    slide_index: int,
    shape_id: int,
    container_type: str,
    table_row: int | None = None,
    table_col: int | None = None,
) -> None:
    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        location_id = ":".join(
            [
                str(slide_index),
                str(shape_id),
                container_type,
                str(table_row),
                str(table_col),
                str(paragraph_index),
            ]
        )
        result[location_id] = paragraph


def _apply_paragraph_decisions(run_items: list[dict[str, Any]], decisions: list[PptxReplacementDecision]) -> int:
    edits_by_run: dict[int, list[tuple[int, int, str]]] = {}
    for decision in decisions:
        candidate = decision.candidate
        for position, run_index in enumerate(candidate.affected_run_indices):
            if run_index >= len(run_items):
                continue
            run_item = run_items[run_index]
            local_start = max(candidate.char_start, run_item["char_start"]) - run_item["char_start"]
            local_end = min(candidate.char_end, run_item["char_end"]) - run_item["char_start"]
            slice_text = decision.replacement if position == 0 else ""
            edits_by_run.setdefault(run_index, []).append((local_start, local_end, slice_text))

    changed = 0
    for run_index, edits in edits_by_run.items():
        run = run_items[run_index]["run"]
        text = run_items[run_index]["text"]
        for start, end, slice_text in sorted(edits, key=lambda item: item[0], reverse=True):
            text = text[:start] + slice_text + text[end:]
        run.text = text
        changed += 1
    return changed


def _paragraph_runs_with_offsets(paragraph: Any) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    cursor = 0
    for run in paragraph.runs:
        text = run.text or ""
        if text:
            items.append({"text": text, "char_start": cursor, "char_end": cursor + len(text), "run": run})
        cursor += len(text)
    return items


def _find_residual_text(
    output_path: Path,
    checked_decisions: list[PptxReplacementDecision],
) -> dict[str, set[str]]:
    """Scan every .xml/.rels part of the saved output package for leftover
    original text from checked_decisions.
    """
    residual: dict[str, set[str]] = {}
    texts_by_category: dict[str, set[str]] = {}
    for decision in checked_decisions:
        text = decision.candidate.text
        if text:
            texts_by_category.setdefault(decision.candidate.category, set()).add(text)
    if not texts_by_category:
        return residual

    with zipfile.ZipFile(output_path) as archive:
        for name in archive.namelist():
            if not name.endswith((".xml", ".rels")):
                continue
            content = archive.read(name).decode("utf-8", errors="ignore")
            found = {
                category
                for category, texts in texts_by_category.items()
                if any(text in content for text in texts)
            }
            if found:
                residual[name] = found
    return residual


def _format_residual_detail(residual: dict[str, set[str]]) -> str:
    return "、".join(
        f"{part_name}: {', '.join(sorted(categories))}"
        for part_name, categories in sorted(residual.items())
    )


_PLACEHOLDER_TYPE_LABELS = {
    "TITLE": "タイトル",
    "CENTER_TITLE": "タイトル",
    "SUBTITLE": "サブタイトル",
    "BODY": "本文",
    "OBJECT": "本文",
}


def pptx_candidate_location_label(candidate: PptxCandidate) -> str:
    if candidate.property_name is not None:
        return f"文書プロパティ({candidate.property_name})"

    slide_label = f"スライド{candidate.slide_index + 1}"
    if candidate.container_type == "notes":
        return f"{slide_label} スピーカーノート 段落{candidate.paragraph_index + 1}"

    if candidate.container_type == "table_cell":
        return (
            f"{slide_label} 表(図形{candidate.shape_id}) "
            f"行{(candidate.table_row or 0) + 1}列{(candidate.table_col or 0) + 1}"
        )

    shape_label = candidate.shape_name
    if candidate.placeholder_type:
        raw_type = candidate.placeholder_type.split(".")[-1] if "." in candidate.placeholder_type else candidate.placeholder_type
        shape_label = _PLACEHOLDER_TYPE_LABELS.get(raw_type, candidate.shape_name)
    return f"{slide_label} {shape_label} 段落{candidate.paragraph_index + 1}"


def pptx_finding_status(decision: PptxReplacementDecision) -> str:
    requires_review = candidate_requires_review(decision.candidate)
    if decision.enabled:
        return "要確認・承認済み" if requires_review else "自動変換"
    if decision.excluded:
        return "確認済み(除外)"
    return "要確認(未処理)" if requires_review else "維持(手動)"


def pptx_finding_reason(decision: PptxReplacementDecision) -> str:
    candidate = decision.candidate
    base = f"{candidate.detection_rule}(信頼度{candidate.confidence:.2f})"
    suffix = {
        "要確認・承認済み": " / 要確認候補を利用者が確認のうえ承認",
        "要確認(未処理)": " / 要確認候補が未承認のため対象外",
        "確認済み(除外)": " / 要確認候補を利用者が確認のうえ除外(原文を維持)",
        "維持(手動)": " / 利用者が対象を解除し原文を維持",
        "自動変換": "",
    }[pptx_finding_status(decision)]
    return base + suffix


def _make_pptx_output_path(source_path: Path, output_dir: Path | None = None, options: ProcessingOptions | None = None) -> Path:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    folder = output_dir if output_dir is not None else source_path.parent
    folder.mkdir(parents=True, exist_ok=True)
    mode_label = (options or ProcessingOptions()).mode_label
    candidate = folder / f"{source_path.stem}_{mode_label}_匿名化_{timestamp}{PPTX_SUPPORTED_EXTENSION}"
    if not candidate.exists():
        return candidate
    return folder / f"{source_path.stem}_{mode_label}_匿名化_{timestamp}_{datetime.now().microsecond:06d}{PPTX_SUPPORTED_EXTENSION}"


def write_pptx_findings_csv(path: Path, decisions: list[PptxReplacementDecision]) -> None:
    headers = ["変換", "処理状態", "位置", "種類", "検査", "検出値", "変換後", "理由"]
    with path.open("w", encoding="utf-8-sig", newline="") as output:
        writer = csv.writer(output)
        writer.writerow(headers)
        for decision in decisions:
            candidate = decision.candidate
            writer.writerow(
                [
                    "対象" if decision.enabled else "維持",
                    pptx_finding_status(decision),
                    pptx_candidate_location_label(candidate),
                    candidate.category,
                    candidate.detection_rule,
                    candidate.text,
                    decision.replacement,
                    pptx_finding_reason(decision),
                ]
            )


def write_pptx_processing_report(
    path: Path,
    *,
    source_path: Path,
    result: PptxConversionResult,
    decisions: list[PptxReplacementDecision],
    inventory: PptxStructureInventory,
    options: ProcessingOptions,
) -> None:
    lines = [
        "PPTX匿名化 処理報告書",
        "",
        f"処理モード: {options.mode_label}",
        f"入力ファイル名: {source_path.name}",
        f"出力ファイル名: {result.output_path.name}",
        f"処理日時: {datetime.now():%Y-%m-%d %H:%M:%S}",
        f"スライド数: {inventory.slide_count}",
        f"検出候補数: {len(decisions)}",
        f"変換run数: {result.converted_run_count}",
        f"変換文書プロパティ数: {result.converted_property_count}",
        f"要確認件数: {result.review_required_count}",
        "",
        "未対応領域:",
    ]
    if inventory.unsupported_features:
        lines.extend(
            f"- {feature.feature_type} (スライド{feature.slide_index + 1}): {feature.count}件"
            for feature in inventory.unsupported_features
        )
    else:
        lines.append("- なし")
    lines.extend(["", "警告内容:"])
    lines.extend([f"- {warning}" for warning in result.warnings] or ["- なし"])
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_pptx_audit_json(
    path: Path,
    *,
    source_path: Path,
    result: PptxConversionResult,
    decisions: list[PptxReplacementDecision],
    inventory: PptxStructureInventory,
    options: ProcessingOptions,
) -> None:
    payload = {
        "schema": PPTX_AUDIT_SCHEMA,
        "summary": {
            "input_file": source_path.name,
            "output_file": result.output_path.name,
            "mode": options.mode_label,
            "processed_at": datetime.now().isoformat(timespec="seconds"),
            "source_sha256": inventory.source_sha256,
            "slide_count": inventory.slide_count,
            "candidate_count": len(decisions),
            "converted_run_count": result.converted_run_count,
            "converted_property_count": result.converted_property_count,
            "review_required_count": result.review_required_count,
            "unsupported_features": [
                {"feature_type": feature.feature_type, "slide_index": feature.slide_index + 1, "count": feature.count}
                for feature in inventory.unsupported_features
            ],
            "warnings": list(result.warnings),
        },
        "findings": [
            {
                "category": decision.candidate.category,
                "location": pptx_candidate_location_label(decision.candidate),
                "detection_rule": decision.candidate.detection_rule,
                "confidence": decision.candidate.confidence,
                "source": decision.candidate.source,
                "status": pptx_finding_status(decision),
                "enabled": decision.enabled,
                "replacement": decision.replacement,
                "reason": pptx_finding_reason(decision),
                "original_hmac_sha256": _pptx_hmac_digest(decision.candidate.text) if decision.candidate.text else "",
            }
            for decision in decisions
        ],
    }
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def validate_pptx_input(path: Path) -> None:
    suffix = path.suffix.lower()
    if suffix == PPTX_MACRO_EXTENSION:
        raise UnsupportedPptxDocumentError(".pptmはPPTX匿名化v1の対象外です。")
    if suffix != PPTX_SUPPORTED_EXTENSION:
        raise UnsupportedPptxDocumentError(f"PPTX匿名化v1は.pptxのみ対応です: {path.name}")
    if not path.exists():
        raise FileNotFoundError(path)


def inspect_pptx_package(path: Path) -> tuple[UnsupportedPptxFeature, ...]:
    features: dict[tuple[str, int], int] = {}
    try:
        presentation = Presentation(path)
        for slide_index, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                _inspect_shape(features, shape, slide_index)
    except Exception as exc:  # pragma: no cover - defensive
        raise UnsupportedPptxDocumentError(f"PPTX構造を読み取れませんでした: {exc}") from exc

    try:
        with zipfile.ZipFile(path) as archive:
            for name in archive.namelist():
                lowered = name.lower()
                if lowered.endswith("vbaproject.bin") or "vbaproject" in lowered:
                    _count_feature(features, "macro", -1)
                if "/comments" in lowered and lowered.endswith(".xml"):
                    _count_feature(features, "comments", -1)
    except zipfile.BadZipFile as exc:
        raise UnsupportedPptxDocumentError(f"PPTX ZIP構造を読み取れませんでした: {exc}") from exc

    return tuple(
        UnsupportedPptxFeature(feature_type=feature_type, slide_index=slide_index, count=count)
        for (feature_type, slide_index), count in sorted(features.items())
    )


def _inspect_shape(features: dict[tuple[str, int], int], shape: Any, slide_index: int) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _inspect_shape(features, child, slide_index)
        return
    if getattr(shape, "has_chart", False):
        _count_feature(features, "chart", slide_index)
        return
    if getattr(shape, "has_table", False) or getattr(shape, "has_text_frame", False):
        return
    if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        _count_feature(features, "embedded_object", slide_index)
        return
    if shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
        _count_feature(features, "smartart", slide_index)
        return
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return
    if shape.shape_type is None:
        # graphicFrame shapes python-pptx can't classify (e.g. some
        # SmartArt/OLE renderings) fall through here.
        _count_feature(features, "unrecognized_shape", slide_index)


def _count_feature(features: dict[tuple[str, int], int], feature_type: str, slide_index: int) -> None:
    key = (feature_type, slide_index)
    features[key] = features.get(key, 0) + 1


def _file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _hmac_digest(value: str) -> str:
    return hmac.new(_state_hmac_key(), value.encode("utf-8"), hashlib.sha256).hexdigest()


def _state_hmac_key() -> bytes:
    env_key = os.environ.get("PRIVACY_CLEANER_STATE_HMAC_KEY")
    if env_key:
        return env_key.encode("utf-8")
    key_path = _state_hmac_key_path()
    if key_path.exists():
        return key_path.read_bytes()
    key = secrets.token_bytes(32)
    try:
        key_path.parent.mkdir(parents=True, exist_ok=True)
        key_path.write_bytes(key)
    except OSError:
        pass
    return key


def _state_hmac_key_path() -> Path:
    local_app_data = os.environ.get("LOCALAPPDATA")
    if local_app_data:
        return Path(local_app_data) / "ExcelPrivacyCleaner" / "state_hmac.key"
    return Path("config") / ".privacy_cleaner_state_hmac.key"


def _pptx_hmac_digest(value: str) -> str:
    return _hmac_digest(f"{PPTX_AUDIT_SCHEMA}:{value}")


def _paragraph_location_id(paragraph: PptxParagraphText) -> str:
    return ":".join(
        [
            str(paragraph.slide_index),
            str(paragraph.shape_id),
            paragraph.container_type,
            str(paragraph.table_row),
            str(paragraph.table_col),
            str(paragraph.paragraph_index),
        ]
    )


def _run_paragraph_location_id(run: PptxTextRun) -> str:
    return ":".join(
        [
            str(run.slide_index),
            str(run.shape_id),
            run.container_type,
            str(run.table_row),
            str(run.table_col),
            str(run.paragraph_index),
        ]
    )


def _runs_by_paragraph_location(runs: tuple[PptxTextRun, ...]) -> dict[str, tuple[PptxTextRun, ...]]:
    grouped: dict[str, list[PptxTextRun]] = {}
    for run in runs:
        grouped.setdefault(_run_paragraph_location_id(run), []).append(run)
    return {key: tuple(value) for key, value in grouped.items()}


def _affected_run_indices(runs: tuple[PptxTextRun, ...], start: int, end: int) -> tuple[int, ...]:
    return tuple(
        run.run_index for run in runs if run.char_start < end and start < run.char_end
    )


def _candidate_id(source_sha256: str, location_id: str, category: str, start: int, end: int, text: str, source: str) -> str:
    payload = "|".join([source_sha256, location_id, category, str(start), str(end), text, source])
    return _hmac_digest(payload)


def _append_text_candidates(
    candidates: list[PptxCandidate],
    seen: set[tuple[str, str, int, int, str, str]],
    inventory: PptxStructureInventory,
    paragraph: PptxParagraphText,
    text: str,
    runs: tuple[PptxTextRun, ...],
    detector: JapanesePresidioDetector,
    ginza_detector: GinzaEntityDetector | None,
) -> None:
    for start, end, category, rule, confidence, original_category in _pptx_detector_results(text, detector, ginza_detector):
        _append_candidate(
            candidates,
            seen,
            inventory,
            paragraph,
            text,
            start,
            end,
            category,
            rule,
            confidence,
            "paragraph",
            affected_runs=_affected_run_indices(runs, start, end),
            original_category=original_category,
        )


def _append_property_candidates(
    candidates: list[PptxCandidate],
    seen: set[tuple[str, str, int, int, str, str]],
    inventory: PptxStructureInventory,
    property_name: str,
    value: str,
    detector: JapanesePresidioDetector,
    ginza_detector: GinzaEntityDetector | None,
) -> None:
    paragraph = PptxParagraphText(
        slide_index=-1,
        shape_id=-1,
        shape_name="",
        container_type="property",
        paragraph_index=0,
        text=value,
    )
    for start, end, category, rule, confidence, original_category in _pptx_detector_results(value, detector, ginza_detector):
        _append_candidate(
            candidates,
            seen,
            inventory,
            paragraph,
            value,
            start,
            end,
            category,
            rule,
            confidence,
            "document_property",
            affected_runs=(),
            original_category=original_category,
            property_name=property_name,
        )


def _append_candidate(
    candidates: list[PptxCandidate],
    seen: set[tuple[str, str, int, int, str, str]],
    inventory: PptxStructureInventory,
    paragraph: PptxParagraphText,
    source_text: str,
    start: int,
    end: int,
    category: str,
    detection_rule: str,
    confidence: float,
    source: str,
    *,
    affected_runs: tuple[int, ...],
    original_category: str,
    property_name: str | None = None,
) -> None:
    if start < 0 or end <= start or end > len(source_text):
        return
    value = source_text[start:end]
    if not value.strip():
        return
    location_id = _paragraph_location_id(paragraph)
    key = (location_id, category, start, end, value, source)
    if key in seen:
        return
    seen.add(key)
    candidates.append(
        PptxCandidate(
            candidate_id=_candidate_id(inventory.source_sha256, location_id, category, start, end, value, source),
            category=category,
            slide_index=paragraph.slide_index,
            shape_id=paragraph.shape_id,
            shape_name=paragraph.shape_name,
            container_type=paragraph.container_type,
            paragraph_index=paragraph.paragraph_index,
            char_start=start,
            char_end=end,
            text=value,
            detection_rule=detection_rule,
            confidence=confidence,
            source=source,
            location_id=location_id,
            affected_run_indices=affected_runs,
            table_row=paragraph.table_row,
            table_col=paragraph.table_col,
            placeholder_type=paragraph.placeholder_type,
            original_category=original_category,
            property_name=property_name,
        )
    )


def _pptx_detector_results(
    text: str,
    detector: JapanesePresidioDetector,
    ginza_detector: GinzaEntityDetector | None = None,
) -> list[tuple[int, int, str, str, float, str]]:
    results: list[tuple[int, int, str, str, float, str]] = []
    for detector_result in detector.analyze(text):
        label = _normalize_pptx_category(entity_label(detector_result.entity_type))
        if label not in PPTX_CANDIDATE_CATEGORIES:
            continue
        start, end = _trim_candidate_span(text, detector_result.start, detector_result.end, label)
        if label == "住所":
            start, end = _extend_address_span(text, start, end)
            if _is_weak_address_candidate(text[start:end]):
                continue
        if _is_generic_pptx_false_positive(label, text[start:end]):
            continue
        results.append((start, end, label, "JapanesePresidioDetector", float(detector_result.score), detector_result.entity_type))
    results.extend(_regex_candidate_results(text))
    if ginza_detector is not None:
        results.extend(_ginza_candidate_results(text, ginza_detector))
    return _select_non_overlapping_pptx_results(text, results)


def _ginza_candidate_results(
    text: str,
    ginza_detector: GinzaEntityDetector,
) -> list[tuple[int, int, str, str, float, str]]:
    results: list[tuple[int, int, str, str, float, str]] = []
    for nlp_result in ginza_detector.analyze(text):
        start, end = _trim_candidate_span(text, nlp_result.start, nlp_result.end, nlp_result.category)
        value = text[start:end]
        if _is_generic_pptx_false_positive(nlp_result.category, value):
            continue
        results.append((start, end, nlp_result.category, WORD_NLP_DETECTION_RULE, WORD_NLP_CONFIDENCE, nlp_result.raw_label))
    return results


def _regex_candidate_results(text: str) -> list[tuple[int, int, str, str, float, str]]:
    results: list[tuple[int, int, str, str, float, str]] = []
    results.extend(_company_designator_results(text))
    results.extend(_contextual_organization_results(text))
    rules = (
        (
            "会社名",
            r"(?:会社名|法人名|取引先|顧客企業)[:：]\s*([一-龯々〆ヵヶぁ-んァ-ヶーA-Za-z0-9]{2,24})",
            "pptx_company_label",
            0.82,
        ),
        (
            "銀行名",
            r"(?:銀行名|金融機関|振込先|口座)[:：は\s]*([一-龯々〆ヵヶぁ-んァ-ヶーA-Za-z0-9]{2,20}(?:銀行|信用金庫|信用組合))",
            "pptx_bank_label",
            0.90,
        ),
        (
            "銀行名",
            r"([一-龯々〆ヵヶぁ-んァ-ヶーA-Za-z0-9]{2,20}(?:銀行|信用金庫|信用組合))",
            "pptx_bank_name",
            0.86,
        ),
        (
            "氏名",
            r"([一-龯々〆ヵヶ]{1,4}[ 　]+[一-龯々〆ヵヶ]{1,5})",
            "pptx_japanese_full_name_space",
            0.70,
        ),
        (
            "氏名",
            r"(?:担当者|連絡者|確認者|責任者|代表者)\s*[:：は]\s*([一-龯々〆ヵヶ]{2,6})(?![ 　])",
            "pptx_person_role_label",
            0.76,
        ),
    )
    for category, pattern, rule, confidence in rules:
        for match in re.finditer(pattern, text):
            start, end = match.span(1) if match.lastindex else match.span()
            start, end = _trim_candidate_span(text, start, end, category)
            value = text[start:end]
            if _is_generic_pptx_false_positive(category, value):
                continue
            if rule == "pptx_japanese_full_name_space" and (
                _is_form_label_before_colon(text, end) or _is_article_number_fragment(text, start)
            ):
                continue
            results.append((start, end, category, rule, confidence, category))
    return results


def _is_form_label_before_colon(text: str, end: int) -> bool:
    limit = min(len(text), end + 15)
    index = end
    while index < limit:
        char = text[index]
        if char in "：:":
            return True
        if char in " \t　" or "一" <= char <= "龯" or char in "々〆ヵヶ":
            index += 1
            continue
        return False
    return False


_ARTICLE_NUMBER_PREFIX_RE = re.compile(r"第[0-9０-９]+$")


def _is_article_number_fragment(text: str, start: int) -> bool:
    prefix = text[max(0, start - 6) : start]
    return bool(_ARTICLE_NUMBER_PREFIX_RE.search(prefix))


def _company_designator_results(text: str) -> list[tuple[int, int, str, str, float, str]]:
    results: list[tuple[int, int, str, str, float, str]] = []
    designator_pattern = "|".join(re.escape(value) for value in COMPANY_DESIGNATORS)
    suffix_pattern = "|".join(re.escape(value) for value in COMPANY_SUFFIX_DESIGNATORS)

    for match in re.finditer(rf"(?:{designator_pattern})[{PPTX_NAME_CHARS}]{{2,24}}", text):
        start, end = _trim_candidate_span(text, match.start(), match.end(), "会社名")
        value = text[start:end]
        if _has_company_name_body(value) and not _is_generic_pptx_false_positive("会社名", value):
            results.append((start, end, "会社名", "pptx_company_prefix", 0.88, "会社名"))

    for match in re.finditer(rf"[{PPTX_NAME_CHARS}]{{2,30}}(?:{suffix_pattern})", text):
        suffix_start = _company_suffix_start(text, match.start(), match.end())
        start, end = _trim_candidate_span(text, suffix_start, match.end(), "会社名")
        value = text[start:end]
        if _has_company_name_body(value) and not _is_generic_pptx_false_positive("会社名", value):
            results.append((start, end, "会社名", "pptx_company_suffix", 0.86, "会社名"))
    return results


def _company_suffix_start(text: str, start: int, end: int) -> int:
    segment = text[start:end]
    adjusted = start
    for index, char in enumerate(segment):
        if char in " \t　:：、。，,/／「」『』（）()はをがにへでと":
            adjusted = start + index + 1
    return adjusted


def _has_company_name_body(value: str) -> bool:
    normalized = value.strip()
    for designator in COMPANY_DESIGNATORS:
        if normalized.startswith(designator):
            return len(normalized) - len(designator) >= 2
        if normalized.endswith(designator):
            return len(normalized) - len(designator) >= 2
    return False


def _contextual_organization_results(text: str) -> list[tuple[int, int, str, str, float, str]]:
    results: list[tuple[int, int, str, str, float, str]] = []
    context_pattern = r"(?:組織として|取引先(?:として|は)?|販売先(?:として|は)?|委託先(?:として|は)?|共同研究先(?:として|は)?)\s*"
    organization_pattern = rf"([{PPTX_NAME_CHARS}]{{2,20}}(?:ラボ|研究所|センター|支店|本社|工場))"
    for match in re.finditer(context_pattern + organization_pattern, text):
        start, end = _trim_candidate_span(text, match.start(1), match.end(1), "会社名")
        value = text[start:end]
        if not any(designator in value for designator in COMPANY_DESIGNATORS) and not _is_generic_pptx_false_positive("会社名", value):
            results.append((start, end, "会社名", "pptx_contextual_organization", 0.72, "会社名"))
    return results


def _select_non_overlapping_pptx_results(
    text: str,
    results: list[tuple[int, int, str, str, float, str]],
) -> list[tuple[int, int, str, str, float, str]]:
    candidates = [
        item
        for item in results
        if 0 <= item[0] < item[1] <= len(text) and len(text[item[0] : item[1]].strip()) >= 2
    ]
    pattern_based = [item for item in candidates if item[3] != WORD_NLP_DETECTION_RULE]
    nlp_based = [item for item in candidates if item[3] == WORD_NLP_DETECTION_RULE]

    pattern_based.sort(key=lambda item: (item[4], item[1] - item[0]), reverse=True)
    selected: list[tuple[int, int, str, str, float, str]] = []
    occupied_by_category: dict[str, list[range]] = {}
    for item in pattern_based:
        start, end, category, *_rest = item
        current = range(start, end)
        if any(current.start < existing.stop and existing.start < current.stop for existing in occupied_by_category.get(category, [])):
            continue
        selected.append(item)
        occupied_by_category.setdefault(category, []).append(current)

    occupied_any = [range(item[0], item[1]) for item in selected]
    nlp_based.sort(key=lambda item: item[1] - item[0], reverse=True)
    for item in nlp_based:
        start, end, *_rest = item
        current = range(start, end)
        if any(current.start < existing.stop and existing.start < current.stop for existing in occupied_any):
            continue
        selected.append(item)
        occupied_any.append(current)

    selected.sort(key=lambda item: (item[0], item[1], item[2]))
    return selected


_LABEL_PREFIX_PATTERNS = {
    "会社名": r"^(?:会社名|法人名|取引先|顧客企業)[:：]\s*",
    "住所": r"^(?:住所|所在地)[:：は]?\s*",
    "銀行名": r"^(?:銀行名|金融機関|振込先|口座)[:：は]?\s*",
}

_SENTENCE_BOUNDARY_RE = re.compile(r"[。、,，.．!！?？\n]")


def _normalize_pptx_category(label: str) -> str:
    mapping = {
        "PERSON": "氏名",
        "個人名": "氏名",
        "組織名": "会社名",
        "電話": "電話番号",
        "メール": "メールアドレス",
        "銀行口座": "銀行名",
        "金融機関": "銀行名",
    }
    return mapping.get(label, label)


def _trim_candidate_span(text: str, start: int, end: int, category: str) -> tuple[int, int]:
    segment = text[start:end]
    prefix_pattern = _LABEL_PREFIX_PATTERNS.get(category)
    if prefix_pattern:
        match = re.match(prefix_pattern, segment)
        if match:
            start += match.end()
            segment = text[start:end]

    if category in {"会社名", "住所", "銀行名"}:
        boundary = _SENTENCE_BOUNDARY_RE.search(segment)
        if boundary:
            end = start + boundary.start()
            segment = text[start:end]

    if category in {"会社名", "氏名"}:
        for honorific in ("様", "さん", "殿", "氏"):
            if segment.endswith(honorific):
                end -= len(honorific)
                segment = text[start:end]
                break

    while start < end and text[start] in " \t　:：、。":
        start += 1
    while end > start and text[end - 1] in " \t　:：、。":
        end -= 1
    return start, end


_INSTITUTION_WORDS = ("銀行", "支店", "会社", "法人", "研究所", "センター", "ラボ", "工場")


def _is_generic_pptx_false_positive(category: str, value: str) -> bool:
    normalized = value.strip()
    if not normalized:
        return True
    if category == "会社名":
        if normalized in COMPANY_DESIGNATORS:
            return True
        if normalized.endswith("担当") or normalized.endswith("先"):
            return len(normalized) <= 3
    if category == "氏名":
        if any(word in normalized for word in _INSTITUTION_WORDS):
            return True
    return False


_ADDRESS_STRENGTH_RE = re.compile(r"[0-9０-９]|丁目|番地|号|〒")
_ADDRESS_SUFFIX_RE = re.compile(r"(ビル|マンション|号館|棟|タワー|ハイツ|レジデンス|階|号室)[0-9０-９A-Za-z一-龯]*")
_WEAK_ADDRESS_WORDS = ("制度", "政策", "説明", "について")


def _extend_address_span(text: str, start: int, end: int) -> tuple[int, int]:
    if not _address_has_strength(text[start:end]):
        return start, end
    match = _ADDRESS_SUFFIX_RE.match(text[end : end + 20])
    if match:
        end += match.end()
    return start, end


def _address_has_strength(value: str) -> bool:
    return bool(_ADDRESS_STRENGTH_RE.search(value))


def _is_weak_address_candidate(value: str) -> bool:
    if not _address_has_strength(value):
        return True
    return any(word in value for word in _WEAK_ADDRESS_WORDS)
