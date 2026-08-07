from __future__ import annotations

import os
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")

from pptx import Presentation  # noqa: E402
from PySide6.QtCore import Qt  # noqa: E402
from PySide6.QtWidgets import QApplication, QMessageBox  # noqa: E402

from excel_privacy_cleaner.pptx_processor import PptxPrivacyProcessor, candidate_requires_review  # noqa: E402
from excel_privacy_cleaner.qt_app import ExcelPrivacyCleanerWindow  # noqa: E402


def assert_true(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def create_review_required_fixture(path: Path) -> None:
    prs = Presentation()
    # python-pptx's default template metadata ("Steve Canny") would
    # otherwise show up as an unrelated 氏名 candidate -- see
    # tests/test_pptx_replacement.py's _new_presentation() for the same note.
    prs.core_properties.last_modified_by = ""
    prs.core_properties.comments = ""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, 4000000, 1000000)
    box.text_frame.text = "場所　アルプススチール様1階応接室"
    # A second, unambiguous high-confidence finding so the slide still has
    # something left to convert after the review-required one is excluded.
    box2 = slide.shapes.add_textbox(0, 1200000, 4000000, 1000000)
    box2.text_frame.text = "連絡先電話は090-1111-2222です。"
    prs.save(path)


def _silence_message_boxes(monkeypatch) -> list[tuple[str, str]]:
    calls: list[tuple[str, str]] = []
    for method in ("information", "warning", "critical", "question"):
        def make_stub(name: str):
            def stub(*_args, **_kwargs):
                calls.append((name, ""))
                return QMessageBox.Yes if name == "question" else None
            return stub
        monkeypatch.setattr(QMessageBox, method, staticmethod(make_stub(method)))
    return calls


def test_pptx_review_required_candidate_exclude_checkbox_resolves_block(monkeypatch) -> None:
    app = QApplication.instance() or QApplication([])
    calls = _silence_message_boxes(monkeypatch)

    with tempfile.TemporaryDirectory(prefix="qt_app_pptx_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_review_required_fixture(source)

        window = ExcelPrivacyCleanerWindow()
        try:
            window.set_source(source)
            assert_true(isinstance(window.processor, PptxPrivacyProcessor), "Selecting a .pptx should create a PptxPrivacyProcessor")

            window.scan_file()
            assert_true(len(window.findings) >= 2, "Both candidates should be detected")
            assert_true(len(window.pptx_decisions) == len(window.findings), "pptx_decisions should track 1:1 with findings")
            review_row = next(
                row
                for row, decision in enumerate(window.pptx_decisions)
                if candidate_requires_review(decision.candidate)
            )
            assert_true(not window.findings[review_row].enabled, "Review-required candidate should default to disabled")

            # Leaving it unresolved must block conversion.
            window.convert_file()
            assert_true(any(name == "critical" for name, _ in calls), "Leaving a review-required candidate unreviewed must block conversion")
            calls.clear()

            # Check "変換しない" for that row -- this is the fix for the block.
            window.table.item(review_row, 1).setCheckState(Qt.Checked)
            assert_true(window.pptx_decisions[review_row].excluded, "Decision should now be marked excluded")
            assert_true(not window.pptx_decisions[review_row].enabled, "Excluded decision must stay disabled")
            assert_true(
                window.table.item(review_row, 0).checkState() == Qt.Unchecked,
                "Checking 変換しない must clear 変換する for the same row",
            )
            assert_true(
                window.table.item(review_row, 5).text() == "確認済み(除外)",
                "Status cell should reflect the excluded state",
            )

            window.convert_file()
            assert_true(any(name == "information" for name, _ in calls), "Conversion should now succeed")
            assert_true(window.history.count() == 1, "Conversion should append one history entry")
        finally:
            window.processor.cleanup()
            window.close()
            app.processEvents()
