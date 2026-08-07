from __future__ import annotations

import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from excel_privacy_cleaner.excel_processor import ProcessingOptions  # noqa: E402
from excel_privacy_cleaner.pptx_processor import (  # noqa: E402
    PptxPrivacyProcessor,
    candidate_requires_review,
    write_pptx_findings_csv,
    write_pptx_processing_report,
    write_pptx_audit_json,
    pptx_candidate_location_label,
    pptx_finding_status,
    _find_residual_text,
)


def assert_true(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def _new_presentation() -> Presentation:
    # python-pptx's default template ships with real incidental metadata
    # (last_modified_by="Steve Canny", comments="generated using
    # python-pptx") that would otherwise show up as an unrelated 氏名
    # candidate in every fixture. Clear it so fixtures only contain the
    # PII they're deliberately built to test.
    prs = Presentation()
    prs.core_properties.last_modified_by = ""
    prs.core_properties.comments = ""
    return prs


def create_replacement_fixture(path: Path) -> None:
    prs = _new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text_frame.text = "株式会社未来会議"

    body = slide.placeholders[1]
    body.text_frame.text = "連絡先電話は090-1111-2222です。"

    # A candidate deliberately split across multiple runs, to exercise the
    # first-run-gets-replacement / later-runs-emptied formatting policy.
    mixed_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
    mixed_paragraph = mixed_box.text_frame.paragraphs[0]
    mixed_paragraph.add_run().text = "山田"
    space_run = mixed_paragraph.add_run()
    space_run.text = " 太郎"
    space_run.font.bold = True
    space_run.font.size = Pt(14)

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(5), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text_frame.text = "会社"
    table_shape.table.cell(0, 1).text_frame.text = "氏名"
    table_shape.table.cell(1, 0).text_frame.text = "株式会社未来会議"
    table_shape.table.cell(1, 1).text_frame.text = "担当者は佐藤花子です。"

    notes = slide.notes_slide
    notes.notes_text_frame.text = "責任者：山田太郎"

    # A second slide reusing the same company/name (as separate lines --
    # the company-designator regex greedily matches up to 24 trailing
    # kanji/hiragana chars with no word-boundary stop, so a single flowing
    # sentence like "株式会社未来会議の担当者は山田太郎です" gets swallowed
    # whole and collides with the adjoining name candidate; Word's own
    # equivalent fixture avoids this the same way, via a run boundary),
    # to exercise alias reuse across slides.
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text_frame.text = "続き"
    slide2.placeholders[1].text_frame.text = "株式会社未来会議\n山田太郎"

    prs.save(path)


def create_bare_company_name_fixture(path: Path) -> None:
    prs = _new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "場所　アルプススチール様1階応接室"
    prs.save(path)


def create_high_confidence_only_fixture(path: Path) -> None:
    prs = _new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "連絡先電話は090-1111-2222です。"
    prs.save(path)


def create_chart_fixture(path: Path) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = _new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("Series 1", (1, 2))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)
    prs.save(path)


def _enable_review_required(decisions) -> list:
    for decision in decisions:
        decision.enabled = True
    return decisions


def test_pptx_replacement_across_shape_table_and_notes() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_replacement_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_replacement_fixture(source)

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        categories = {decision.candidate.original_category or decision.candidate.category for decision in decisions}
        assert_true(any(decision.candidate.text == "090-1111-2222" for decision in decisions), "Phone number should be detected")
        assert_true(any(decision.candidate.text == "株式会社未来会議" for decision in decisions), "Company name should be detected")
        assert_true(any("山田" in decision.candidate.text or "太郎" in decision.candidate.text for decision in decisions), "Name should be detected")

        name_replacements = {
            decision.replacement
            for decision in decisions
            if "山田" in decision.candidate.text and "太郎" in decision.candidate.text
        }
        company_replacements = {
            decision.replacement for decision in decisions if decision.candidate.text == "株式会社未来会議"
        }
        assert_true(len(name_replacements) <= 1, "Same person across slides should reuse the same alias")
        assert_true(len(company_replacements) == 1, "Same company across slides/table should reuse the same alias")

        _enable_review_required(decisions)
        result = processor.convert(source, decisions, output_dir=tmp)
        assert_true(result.output_path.exists(), "Anonymized pptx should be written")
        assert_true(result.csv_path.exists(), "CSV should be written")
        assert_true(result.report_path.exists(), "Report should be written")
        assert_true(result.warnings == (), f"Should have no warnings, got: {result.warnings}")

        reopened = Presentation(result.output_path)
        slide1 = reopened.slides[0]
        all_text = []
        for shape in slide1.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        all_text.append(cell.text_frame.text)
        assert_true(slide1.has_notes_slide, "Notes slide should still exist")
        all_text.append(slide1.notes_slide.notes_text_frame.text)
        combined = "\n".join(all_text)

        for original in ("株式会社未来会議", "090-1111-2222", "山田", "太郎", "佐藤花子"):
            assert_true(original not in combined, f"Original text should not remain: {original}")


def test_split_run_replacement_preserves_sibling_formatting() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_run_format_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_replacement_fixture(source)

        before = Presentation(source)
        textbox = [
            shape
            for shape in before.slides[0].shapes
            if shape.has_text_frame and shape.text_frame.text == "山田 太郎"
        ][0]
        bold_run_before = textbox.text_frame.paragraphs[0].runs[1]
        assert_true(bold_run_before.font.bold is True, "Fixture sanity check: second run should start bold")

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        _enable_review_required(decisions)
        result = processor.convert(source, decisions, output_dir=tmp)

        reopened = Presentation(result.output_path)
        textbox_after = [
            shape
            for shape in reopened.slides[0].shapes
            if shape.has_text_frame and "個人" in shape.text_frame.text and shape.shape_id == textbox.shape_id
        ]
        assert_true(len(textbox_after) == 1, "Textbox should still be present after conversion")
        bold_run_after = textbox_after[0].text_frame.paragraphs[0].runs[1]
        assert_true(bold_run_after.font.bold is True, "Sibling run formatting (bold) should be preserved")
        assert_true(bold_run_after.font.size == Pt(14), "Sibling run formatting (size) should be preserved")


def test_low_confidence_candidate_blocks_by_default_and_review_required_allows_conversion() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_low_confidence_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_bare_company_name_fixture(source)

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        assert_true(len(decisions) > 0, "Bare company name should still produce a low-confidence candidate")
        assert_true(all(not decision.enabled for decision in decisions), "Low-confidence candidates should start disabled")
        assert_true(
            all(candidate_requires_review(decision.candidate) for decision in decisions),
            "Candidates in this fixture should all be below the review threshold",
        )

        try:
            processor.convert(source, decisions, output_dir=tmp)
        except RuntimeError as exc:
            assert_true("未確認候補" in str(exc), "Should block on unresolved review-required candidates")
        else:
            raise AssertionError("Conversion should be blocked while review-required candidates are unresolved")

        # A second scan (since convert() didn't consume the temp copy) to
        # approve and confirm conversion succeeds once reviewed.
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        _enable_review_required(decisions)
        result = processor.convert(source, decisions, output_dir=tmp)
        assert_true(result.output_path.exists(), "Conversion should succeed once review-required candidates are approved")


def test_excluded_candidate_keeps_original_text() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_excluded_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_bare_company_name_fixture(source)

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        for decision in decisions:
            decision.enabled = False
            decision.excluded = True

        result = processor.convert(source, decisions, output_dir=tmp)
        reopened = Presentation(result.output_path)
        combined = "\n".join(
            shape.text_frame.text for shape in reopened.slides[0].shapes if shape.has_text_frame
        )
        assert_true("アルプススチール" in combined, "Excluded candidate should keep its original text")


def test_high_confidence_candidate_auto_converts_without_review() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_high_confidence_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_high_confidence_only_fixture(source)

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        assert_true(all(decision.enabled for decision in decisions), "High-confidence phone match should auto-enable")
        result = processor.convert(source, decisions, output_dir=tmp)
        reopened = Presentation(result.output_path)
        combined = "\n".join(
            shape.text_frame.text for shape in reopened.slides[0].shapes if shape.has_text_frame
        )
        assert_true("090-1111-2222" not in combined, "Phone number should be converted")


def test_unsupported_chart_blocks_external_share_and_warns_analysis() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_unsupported_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_chart_fixture(source)

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        assert_true(len(processor.inventory.unsupported_features) > 0, "Chart should be recorded as unsupported")

        result = processor.convert(source, decisions, output_dir=tmp)
        assert_true(
            any("未対応領域" in warning for warning in result.warnings),
            "Analysis mode should warn about unsupported regions",
        )

        processor2 = PptxPrivacyProcessor()
        decisions2 = processor2.scan(source, options=ProcessingOptions(mode="external"))
        try:
            processor2.convert(source, decisions2, output_dir=tmp)
        except RuntimeError as exc:
            assert_true("外部共有用" in str(exc), "External-share mode should block on unsupported regions")
        else:
            raise AssertionError("External-share conversion should be blocked by the chart")


def test_residual_text_detection() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_residual_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_high_confidence_only_fixture(source)

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        result = processor.convert(source, decisions, output_dir=tmp)

        residual = _find_residual_text(result.output_path, decisions)
        assert_true(residual == {}, f"No residual text should remain, got: {residual}")

        # Inject the original phone number back into a slide XML part to
        # prove the residual scan actually catches leftover text.
        tmp_path = result.output_path.with_name(result.output_path.name + ".tmp")
        with zipfile.ZipFile(result.output_path) as source_archive, zipfile.ZipFile(
            tmp_path, "w", compression=zipfile.ZIP_DEFLATED
        ) as target_archive:
            for info in source_archive.infolist():
                data = source_archive.read(info.filename)
                if info.filename == "docProps/app.xml":
                    data = data.replace(b"</Properties>", b"<Notes>090-1111-2222</Notes></Properties>")
                target_archive.writestr(info, data)
        tmp_path.replace(result.output_path)

        residual_after_injection = _find_residual_text(result.output_path, decisions)
        assert_true(residual_after_injection != {}, "Residual scan should detect the injected leftover text")


def test_csv_report_and_audit_json_output() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_artifacts_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_high_confidence_only_fixture(source)

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        result = processor.convert(source, decisions, output_dir=tmp)

        csv_text = result.csv_path.read_text(encoding="utf-8-sig")
        assert_true("変換" in csv_text and "検査" in csv_text, "CSV should have the expected headers")
        assert_true("090-1111-2222" in csv_text, "CSV should list the detected original value")

        report_text = result.report_path.read_text(encoding="utf-8")
        assert_true("PPTX匿名化 処理報告書" in report_text, "Report should identify itself")

        audit_path = result.output_path.with_name(f"{result.output_path.stem}_監査報告.json")
        import json

        payload = json.loads(audit_path.read_text(encoding="utf-8"))
        assert_true(payload["schema"] == "pptx_audit_v1", "Audit JSON should carry the pptx schema tag")
        assert_true(
            all("090-1111-2222" not in json.dumps(finding) for finding in payload["findings"]),
            "Audit JSON must not contain raw original text",
        )


def test_pptx_candidate_location_label_and_status_helpers() -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_helpers_") as tmpdir:
        tmp = Path(tmpdir)
        source = tmp / "fixture.pptx"
        create_replacement_fixture(source)

        processor = PptxPrivacyProcessor()
        decisions = processor.scan(source, options=ProcessingOptions(mode="analysis"))
        for decision in decisions:
            label = pptx_candidate_location_label(decision.candidate)
            assert_true("スライド" in label or "文書プロパティ" in label, f"Location label should be human-readable: {label}")
            status = pptx_finding_status(decision)
            assert_true(status in {"自動変換", "要確認(未処理)", "要確認・承認済み", "確認済み(除外)", "維持(手動)"}, f"Unexpected status: {status}")


if __name__ == "__main__":
    import pytest

    raise SystemExit(pytest.main([__file__, "-v"]))
